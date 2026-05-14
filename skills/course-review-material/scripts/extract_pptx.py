"""Extract content from .pptx files to structured Markdown.

Extracts text from all shapes, exports embedded images, and renders
slides as PNG screenshots for multimodal recognition.

Usage:
    python extract_pptx.py <input.pptx> [--output-dir <dir>] [--render-slides] [--dpi <dpi>]
"""

import argparse
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def extract_slide_text_and_images(slide, slide_num, images_dir):
    slide_text = []
    image_refs = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.replace("\n", " ") for cell in row.cells]
                rows.append(cells)
            if rows:
                slide_text.append("\n| " + " | ".join(rows[0]) + " |")
                slide_text.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                for row in rows[1:]:
                    slide_text.append("| " + " | ".join(row) + " |")
                slide_text.append("")
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_text.append(text)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = os.path.splitext(image.content_type.split("/")[-1])[1]
            if not ext:
                ext = ".png"
            img_name = f"slide{slide_num:02d}_img_{len(image_refs) + 1:03d}{ext}"
            img_path = os.path.join(images_dir, img_name)
            with open(img_path, "wb") as f:
                f.write(image.blob)
            image_refs.append(img_name)

    return slide_text, image_refs


def extract_pptx(filepath, output_dir, render_slides=False, dpi=150):
    prs = Presentation(filepath)

    os.makedirs(output_dir, exist_ok=True)
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    if render_slides:
        slides_dir = os.path.join(output_dir, "slides")
        os.makedirs(slides_dir, exist_ok=True)
        _render_slides_to_png(filepath, slides_dir, dpi)

    md_lines = []
    raw_texts = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_header = f"\n## Slide {idx}\n"
        md_lines.append(slide_header)

        slide_text, image_refs = extract_slide_text_and_images(slide, idx, images_dir)

        if slide_text:
            for line in slide_text:
                md_lines.append(line + "\n")
                raw_texts.append(line)

        for img_ref in image_refs:
            md_lines.append(f"![Slide {idx} image](images/{img_ref})\n")
            md_lines.append(f"<!-- image source: {img_ref} -->\n")

        if not slide_text and not image_refs and render_slides:
            md_lines.append(f"![Slide {idx} screenshot](slides/slide_{idx:03d}.png)\n")
            md_lines.append(f"<!-- slide source: slide_{idx:03d}.png (needs multimodal recognition) -->\n")

        if not slide_text and not image_refs and not render_slides:
            md_lines.append("*(This slide appears to be image-only. Use --render-slides to enable recognition.)*\n")

    text_dir = os.path.join(output_dir, "text")
    os.makedirs(text_dir, exist_ok=True)
    raw_text_path = os.path.join(text_dir, "raw_text.txt")
    with open(raw_text_path, "w", encoding="utf-8") as f:
        f.write("\n".join(raw_texts))

    extracted_path = os.path.join(output_dir, "extracted.md")
    with open(extracted_path, "w", encoding="utf-8") as f:
        f.write("".join(md_lines))

    print(f"[done] slides: {len(prs.slides)}")
    print(f"[done] raw text: {raw_text_path}")
    print(f"[done] extracted: {extracted_path}")
    return extracted_path


def _render_slides_to_png(pptx_path, slides_dir, dpi):
    try:
        import tempfile

        try:
            from pptx import Presentation as Prs
        except ImportError:
            return

        try:
            import fitz
        except ImportError:
            print("[warn] PyMuPDF (fitz) not installed. Cannot render slides.")
            print("[warn] Install with: pip install PyMuPDF")
            print("[warn] Alternative: export PPTX to PDF manually, then use extract_pdf.py")
            return

        print("[render] Converting PPTX to PDF for slide rendering...")

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_pdf = tmp.name

        import platform
        if platform.system() == "Windows":
            import comtypes.client

            powerpoint = None
            try:
                powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                powerpoint.Visible = False
                deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
                deck.SaveAs(os.path.abspath(tmp_pdf), 32)
                deck.Close()
            finally:
                if powerpoint is not None:
                    try:
                        powerpoint.Quit()
                    except Exception:
                        pass
        else:
            import subprocess

            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", os.path.dirname(tmp_pdf), pptx_path],
                capture_output=True, text=True
            )
            if result.returncode != 0:
                print(f"[warn] LibreOffice conversion failed: {result.stderr}")
                print("[warn] Slide rendering skipped. Install LibreOffice or use --no-render.")
                os.unlink(tmp_pdf)
                return

        doc = fitz.open(tmp_pdf)
        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            slide_path = os.path.join(slides_dir, f"slide_{page_num + 1:03d}.png")
            pix.save(slide_path)

        doc.close()
        os.unlink(tmp_pdf)
        print(f"[render] {len(os.listdir(slides_dir))} slides rendered to {slides_dir}")

    except Exception as e:
        print(f"[warn] Slide rendering failed: {e}")
        print("[warn] Slides will be referenced but not rendered.")


def main():
    parser = argparse.ArgumentParser(description="Extract pptx content to structured Markdown")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("--output-dir", "-o", default=None, help="Output directory (default: {input_name}/)")
    parser.add_argument("--render-slides", "-r", action="store_true", help="Render slides as PNG images")
    parser.add_argument("--dpi", type=int, default=150, help="DPI for slide rendering (default: 150)")
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"Error: file not found: {input_path}")
        sys.exit(1)

    output_dir = args.output_dir or str(input_path.parent / input_path.stem)

    try:
        extract_pptx(str(input_path), output_dir, render_slides=args.render_slides, dpi=args.dpi)
    except Exception as e:
        print(f"Error: failed to process {input_path}: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
